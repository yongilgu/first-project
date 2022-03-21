
from pptx import Presentation
import pandas



df= pandas.read_excel('//10.204.115.19/신소재개발실/1. 분석지원/1. 분석/1. 분석 보고서/2022년 분석업무현황(ver_4).xlsm',sheet_name='업무현황')
prs = Presentation('//10.204.115.19/신소재개발실/1. 분석지원/1. 분석/1. 분석 보고서/work automation prj/auto report_r1.pptm')

#1페이지(표지)
add_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(add_slide_layout)

slide.placeholders[0].text = df.제목[0] 
 #날짜 자동생성 나중에 만들어보자

#2페이지(분석개요)
add_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(add_slide_layout)

slide.placeholders[0].text = "분석개요" 
slide.placeholders[1].text = df.제목[0]
slide.placeholders[10].text = df.의뢰번호[0]
slide.placeholders[11].text = df.분석종류[0]
slide.placeholders[12].text = df.분석구분[0]
slide.placeholders[13].text = df.소재[0]
slide.placeholders[15].text = df.소속[0]
slide.placeholders[16].text = df.신청인[0]
slide.placeholders[17].text = df.고객명[0]
slide.placeholders[18].text = df.제품명[0]
slide.placeholders[20].text = df.장비명[0]
slide.placeholders[21].text = df.공정온도[0]
slide.placeholders[22].text = df.공정명[0]
slide.placeholders[23].text = df.세부공정[0]



#3페이지(종합)
add_slide_layout = prs.slide_layouts[2]
slide = prs.slides.add_slide(add_slide_layout)

slide.placeholders[0].text = "결론" 

#4페이지()
add_slide_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(add_slide_layout)

slide.placeholders[0].text = "형상 분석 결과" 

prs.save('//10.204.115.19/신소재개발실/1. 분석지원/1. 분석/1. 분석 보고서/22-000_소재_종합_.pptm')

#pyinstaller --onefile --noconsole 220218_ver3.py

# [layout0 의Placeholder ]
# 0:Title 1 제목
# 1:Text Placeholder 날짜
# 15:Text Placeholder 팀명

# [layout1 의Placeholder ]
# 0:Title 제목
# 1:Text Placeholder 의뢰내용
# 10:Text Placeholder 의뢰번호
# 11:Text Placeholder 분석종류
# 12:Text Placeholder 분석구분
# 13:Text Placeholder 소재
# 14:Text Placeholder 기타
# 15:Text Placeholder 소속
# 16:Text Placeholder 신청인
# 17:Text Placeholder 고객명
# 18:Text Placeholder 제품명
# 20:Text Placeholder 장비명
# 21:Text Placeholder 공정온도
# 22:Text Placeholder 공정명
# 23:Text Placeholder 세부공정

#[반복적으로 실행되도록 하는 구문이다.]
# for i in df.index: 
#     add_slide_layout = prs.slide_layouts[1]
#     slide = prs.slides.add_slide(add_slide_layout)

#     slide.placeholders[0].text = df.분석구분[i] #번호를 i로 교체해주어야 한다. 만약 특정 행을 가져오려면 0(제목)부터 진행
#     slide.placeholders[1].text = df.제목[i]

#[이건 글자를 입력할떄 사용]
# slide.placeholders[0].text = "분석개요" 
# slide.placeholders[1].text = "실리콘 분석 개요에 대한 개요개요"

# [슬라이드의 주소를 알아내는 용도이다.]
# shapes = slide.shapes 
# for shape in shapes:
#     print(str(shape.placeholder_format.idx)+ ":"+ shape.name)

# 220217 완료!!!!